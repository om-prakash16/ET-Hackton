import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize presentation (Widescreen 16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Define Premium Theme Colors
BG_COLOR = RGBColor(15, 23, 42)       # Slate 900
ACCENT_CYAN = RGBColor(6, 182, 212)   # Cyan 500
ACCENT_GREEN = RGBColor(16, 185, 129) # Emerald 500
ACCENT_RED = RGBColor(239, 68, 68)    # Red 500
TEXT_WHITE = RGBColor(248, 250, 252)  # Slate 50
TEXT_GRAY = RGBColor(148, 163, 184)   # Slate 400

def set_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_title(slide, text, color=ACCENT_CYAN):
    title_shape = slide.shapes.title
    title_shape.text = text
    title_shape.text_frame.paragraphs[0].font.color.rgb = color
    title_shape.text_frame.paragraphs[0].font.name = "Arial"
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    return title_shape

def add_bullet(tf, text, level=0, bold=False, color=TEXT_GRAY, size=24):
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.size = Pt(size)
    p.font.name = "Arial"
    return p

# --- SLIDE 1: Title Slide ---
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_background(slide)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "IndusBrain"
title.text_frame.paragraphs[0].font.color.rgb = ACCENT_CYAN
title.text_frame.paragraphs[0].font.size = Pt(88)
title.text_frame.paragraphs[0].font.bold = True

subtitle.text = "The Multi-Tenant Industrial AI Operating System\nFrom Data Silos to Decision Intelligence"
subtitle.text_frame.paragraphs[0].font.color.rgb = TEXT_WHITE
subtitle.text_frame.paragraphs[0].font.size = Pt(28)
if len(subtitle.text_frame.paragraphs) > 1:
    subtitle.text_frame.paragraphs[1].font.color.rgb = ACCENT_GREEN
    subtitle.text_frame.paragraphs[1].font.size = Pt(22)
    subtitle.text_frame.paragraphs[1].font.italic = True

# --- SLIDE 2: The $100B Problem ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_background(slide)
add_title(slide, "1. The Heavy Industry Problem", ACCENT_RED)
tf = slide.placeholders[1].text_frame

add_bullet(tf, "The Reality of Industrial Data Silos", bold=True, color=TEXT_WHITE, size=28)
add_bullet(tf, "Engineers spend 35% of their time searching for information (McKinsey 2024).", level=1, size=22)
add_bullet(tf, "P&ID Drawings are in AutoCAD, Maintenance Logs in SAP, and OEM Manuals in 400-page PDFs.", level=1, size=22)
add_bullet(tf, "The Consequence of Missing Context", bold=True, color=TEXT_WHITE, size=28)
add_bullet(tf, "Unplanned downtime costs heavy industry $50B+ annually.", level=1, size=22)
add_bullet(tf, "When a pump vibrates, engineers lose hours piecing together scattered data before they can act.", level=1, size=22)

# --- SLIDE 3: The Solution ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_background(slide)
add_title(slide, "2. Introducing IndusBrain OS", ACCENT_CYAN)
tf = slide.placeholders[1].text_frame

add_bullet(tf, "A Unified Command Center", bold=True, color=TEXT_WHITE, size=28)
add_bullet(tf, "IndusBrain is an end-to-end AI platform that ingests unstructured industrial data and turns it into actionable intelligence.", level=1, size=22)
add_bullet(tf, "Universal Document Intelligence Pipeline", bold=True, color=ACCENT_CYAN, size=28)
add_bullet(tf, "Automatically runs OCR & NLP on PDFs, DWGs, and Excel sheets.", level=1, size=22)
add_bullet(tf, "Active Knowledge Graph", bold=True, color=ACCENT_CYAN, size=28)
add_bullet(tf, "Extracts entities to build a topological digital twin (e.g., Pump A connects to Tank B).", level=1, size=22)

# --- SLIDE 4: Standard RAG vs Agentic GraphRAG ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_background(slide)
add_title(slide, "3. The Technical Differentiator", ACCENT_GREEN)
tf = slide.placeholders[1].text_frame

add_bullet(tf, "Standard RAG (The Old Way) fails in engineering.", bold=True, color=ACCENT_RED, size=28)
add_bullet(tf, "Relies purely on vector similarity (keyword matching).", level=1, size=22)
add_bullet(tf, "Cannot answer topological questions like 'What happens downstream if Valve V-200 fails?'", level=1, size=22)
add_bullet(tf, "Agentic GraphRAG (Our Breakthrough)", bold=True, color=ACCENT_GREEN, size=28)
add_bullet(tf, "Combines Vector Embeddings (Qdrant) with Graph Networks (Neo4j).", level=1, size=22)
add_bullet(tf, "The AI acts autonomously, navigating the physical topology of the factory.", level=1, size=22)
add_bullet(tf, "Zero Hallucination: 100% strict citations to exact source documents.", level=1, size=22)

# --- SLIDE 5: Multi-Tenancy & Security ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_background(slide)
add_title(slide, "4. Enterprise Multi-Tenancy & Security", ACCENT_CYAN)
tf = slide.placeholders[1].text_frame

add_bullet(tf, "Built for B2B SaaS Scale", bold=True, color=TEXT_WHITE, size=28)
add_bullet(tf, "Super Admins can instantly provision isolated 'Enterprise Brains' for different clients.", level=1, size=22)
add_bullet(tf, "Strict Data Isolation", bold=True, color=TEXT_WHITE, size=28)
add_bullet(tf, "Client A's Knowledge Graph and Vectors never leak into Client B's AI context.", level=1, size=22)
add_bullet(tf, "Granular Role-Based Access Control (RBAC)", bold=True, color=TEXT_WHITE, size=28)
add_bullet(tf, "JWT-secured architecture enforces permissions.", level=1, size=22)
add_bullet(tf, "E.g., A 'Maintenance Engineer' can upload evidence, while a 'Viewer' is read-only.", level=1, size=22)

# --- SLIDE 6: Automated Root Cause Analysis ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_background(slide)
add_title(slide, "5. Core Feature: Automated RCA", ACCENT_GREEN)
tf = slide.placeholders[1].text_frame

add_bullet(tf, "Turning 3 Days into 3 Seconds", bold=True, color=TEXT_WHITE, size=28)
add_bullet(tf, "When an anomaly occurs (e.g., 15% pressure drop), the AI takes over.", level=1, size=22)
add_bullet(tf, "The AI orchestrator cross-references live telemetry with the OEM Manual and Maintenance Logs.", level=1, size=22)
add_bullet(tf, "The Output: An Instant RCA Tree", bold=True, color=TEXT_WHITE, size=28)
add_bullet(tf, "1. Symptom Identified.", level=1, size=22)
add_bullet(tf, "2. Probable Cause Extracted.", level=1, size=22)
add_bullet(tf, "3. Root Cause Determined & Work Order Recommended.", level=1, size=22)

# --- SLIDE 7: The Architecture Stack ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_background(slide)
add_title(slide, "6. Tech Stack & Architecture", ACCENT_CYAN)
tf = slide.placeholders[1].text_frame

add_bullet(tf, "Frontend Presentation", bold=True, color=TEXT_WHITE, size=26)
add_bullet(tf, "Next.js 14, React, Tailwind CSS, Framer Motion (Highly Responsive UI).", level=1, size=20)
add_bullet(tf, "Backend API & Data Processing", bold=True, color=TEXT_WHITE, size=26)
add_bullet(tf, "FastAPI (Python) for asynchronous AI routing and document ingestion.", level=1, size=20)
add_bullet(tf, "The AI Intelligence Layer", bold=True, color=TEXT_WHITE, size=26)
add_bullet(tf, "Google Gemini 2.5 Flash as the core Agentic Reasoning Engine.", level=1, size=20)
add_bullet(tf, "Databases", bold=True, color=TEXT_WHITE, size=26)
add_bullet(tf, "Neo4j for Graph Topology | Qdrant for Vector Embeddings.", level=1, size=20)

# --- SLIDE 8: Product Roadmap (Flutter) ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_background(slide)
add_title(slide, "7. Product Roadmap: Flutter Mobile App", ACCENT_CYAN)
tf = slide.placeholders[1].text_frame

add_bullet(tf, "Empowering the 'Deskless' Worker", bold=True, color=TEXT_WHITE, size=28)
add_bullet(tf, "Engineers on the factory floor don't use laptops.", level=1, size=22)
add_bullet(tf, "Phase 2 introduces a cross-platform (iOS/Android) mobile interface.", level=1, size=22)
add_bullet(tf, "Scan QR Codes on Machinery", bold=True, color=ACCENT_CYAN, size=28)
add_bullet(tf, "Scan a pump to instantly load its context into the Agentic AI Copilot.", level=1, size=22)
add_bullet(tf, "Upload Photographic Evidence", bold=True, color=ACCENT_GREEN, size=28)
add_bullet(tf, "Snap a picture of a leaking valve to trigger an automated RCA workflow instantly.", level=1, size=22)

# --- SLIDE 9: Business Value & ROI ---
slide = prs.slides.add_slide(prs.slide_layouts[1])
set_background(slide)
add_title(slide, "8. Business Value Proposition", ACCENT_GREEN)
tf = slide.placeholders[1].text_frame

add_bullet(tf, "Massive ROI for Heavy Industry", bold=True, color=TEXT_WHITE, size=28)
add_bullet(tf, "-90% Reduction in Troubleshooting Time", bold=True, color=ACCENT_CYAN, size=24)
add_bullet(tf, "Engineers stop searching and start fixing.", level=1, size=20)
add_bullet(tf, "100% Automated Regulatory Compliance", bold=True, color=ACCENT_CYAN, size=24)
add_bullet(tf, "Graph AI actively audits operations against safety frameworks like OISD-105.", level=1, size=20)
add_bullet(tf, "Prevention of Catastrophic Failure", bold=True, color=ACCENT_CYAN, size=24)
add_bullet(tf, "Proactively links anomalies to historical maintenance data before a machine breaks.", level=1, size=20)

# --- SLIDE 10: Conclusion ---
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_background(slide)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Thank You"
title.text_frame.paragraphs[0].font.color.rgb = ACCENT_GREEN
title.text_frame.paragraphs[0].font.size = Pt(72)
title.text_frame.paragraphs[0].font.bold = True

subtitle.text = "IndusBrain: We didn't just build a chatbot.\nWe built the Industrial OS of the future."
subtitle.text_frame.paragraphs[0].font.color.rgb = TEXT_WHITE
subtitle.text_frame.paragraphs[0].font.size = Pt(28)
if len(subtitle.text_frame.paragraphs) > 1:
    subtitle.text_frame.paragraphs[1].font.color.rgb = ACCENT_CYAN
    subtitle.text_frame.paragraphs[1].font.size = Pt(24)

prs.save('IndusBrain_PitchDeck_Premium.pptx')
print("Successfully generated IndusBrain_PitchDeck_Premium.pptx")
